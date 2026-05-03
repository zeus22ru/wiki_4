# -*- coding: utf-8 -*-
"""
Генерация PowerPoint-версии презентации БочкарИИ (эквивалент index.html).
Запуск из корня репозитория или из этой папки:

    python docs/presentation/generate_pptx.py

Требуется: pip install python-pptx
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

# Пути
HERE = Path(__file__).resolve().parent
DOCS = HERE.parent
IMAGES = DOCS / "images"
OUT = HERE / "BochkarII_presentation.pptx"

# Фирменные цвета (как в styles.css)
CLR_PRIMARY = RGBColor(0x7C, 0x5A, 0x22)
CLR_TEXT = RGBColor(0x25, 0x21, 0x1D)
CLR_MUTED = RGBColor(0x5E, 0x55, 0x4A)
CLR_BG = RGBColor(0xF4, 0xF1, 0xEB)
# Как в styles.css: .bk-surface, .bk-border, .bk-eyebrow, .bk-stat__value / __label
CLR_SURFACE = RGBColor(0xFF, 0xFD, 0xF8)
CLR_BORDER = RGBColor(0xD8, 0xCD, 0xB8)
CLR_TEXT_SOFT = RGBColor(0x5E, 0x55, 0x4A)
CLR_PRIMARY_DEEP = RGBColor(0x5D, 0x42, 0x18)
# rgba(165,130,80,0.14) + border — близкие сплошные оттенки
CLR_EYEBROW_BG = RGBColor(0xED, 0xE6, 0xDA)
CLR_EYEBROW_BORDER = RGBColor(0xC9, 0xB8, 0x9A)
CLR_PRIMARY_SOFT = RGBColor(0xA5, 0x82, 0x50)


def _set_run_font(run, size_pt: float, bold: bool = False, color: RGBColor | None = None):
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    run.font.name = "Calibri"


def _fill_body_placeholder(slide, text: str, left, top, width, height, size_pt: float = 14):
    """Текстовое поле с переносами."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    _set_run_font(run, size_pt, color=CLR_TEXT)
    return box


def _add_title_slide(prs: Presentation):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    # Фон
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = CLR_BG

    # Заголовок
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "БочкарИИ"
    _set_run_font(run, 44, bold=True, color=CLR_PRIMARY)

    sub = slide.shapes.add_textbox(Inches(1.0), Inches(3.2), Inches(11.0), Inches(1.4))
    stf = sub.text_frame
    stf.word_wrap = True
    sp = stf.paragraphs[0]
    sp.alignment = PP_ALIGN.CENTER
    sr = sp.add_run()
    sr.text = (
        "AI-ассистент по корпоративной базе знаний.\n"
        "Один чат вместо десятка вкладок wiki."
    )
    _set_run_font(sr, 20, color=CLR_TEXT)

    meta = slide.shapes.add_textbox(Inches(1.0), Inches(5.0), Inches(11.0), Inches(1.0))
    mtf = meta.text_frame
    mp = mtf.paragraphs[0]
    mp.alignment = PP_ALIGN.CENTER
    mr = mp.add_run()
    mr.text = "Демо для команды  ·  Май 2026  ·  Локальное развёртывание"
    _set_run_font(mr, 14, color=CLR_MUTED)


def _section_header(slide, eyebrow: str, title: str):
    eb = slide.shapes.add_textbox(Inches(0.6), Inches(0.45), Inches(12.0), Inches(0.4))
    ep = eb.text_frame.paragraphs[0]
    er = ep.add_run()
    er.text = eyebrow.upper()
    _set_run_font(er, 11, bold=True, color=CLR_PRIMARY)

    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.75), Inches(12.0), Inches(0.9))
    tp = tb.text_frame.paragraphs[0]
    tr = tp.add_run()
    tr.text = title
    _set_run_font(tr, 28, bold=True, color=CLR_TEXT)


def _blank_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = CLR_BG
    return slide


def _add_eyebrow_pill(slide, text: str, left, top, width, height) -> None:
    """Скруглённый бейдж в стиле .bk-eyebrow."""
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    pill.fill.solid()
    pill.fill.fore_color.rgb = CLR_EYEBROW_BG
    pill.line.color.rgb = CLR_EYEBROW_BORDER
    pill.line.width = Pt(0.75)
    try:
        pill.adjustments[0] = 0.35
    except (AttributeError, IndexError, TypeError):
        pass
    tf = pill.text_frame
    tf.clear()
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text.upper()
    _set_run_font(r, 10, bold=True, color=CLR_PRIMARY)
    try:
        r.font.spacing = Pt(0.35)
    except (AttributeError, TypeError):
        pass


def _add_stat_card(slide, left, top, width, height, value: str, label: str) -> None:
    """Карточка метрики: .bk-stat (белый фон, рамка, скругление; тень в OOXML здесь не задаём)."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = CLR_SURFACE
    card.line.color.rgb = CLR_BORDER
    card.line.width = Pt(0.75)
    try:
        card.adjustments[0] = 0.12
    except (AttributeError, IndexError, TypeError):
        pass
    tf = card.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.14)
    tf.margin_right = Inches(0.14)
    tf.margin_top = Inches(0.18)
    tf.margin_bottom = Inches(0.18)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    p1.space_after = Pt(4)
    rv = p1.add_run()
    rv.text = value
    _set_run_font(rv, 30, bold=True, color=CLR_PRIMARY_DEEP)

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    rl = p2.add_run()
    rl.text = label
    _set_run_font(rl, 11, color=CLR_TEXT_SOFT)


def _flow_node(slide, left, top, width, height, text: str, font_pt: float = 8) -> None:
    """Прямоугольник узла блок-схемы."""
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = CLR_SURFACE
    sh.line.color.rgb = CLR_BORDER
    sh.line.width = Pt(0.75)
    try:
        sh.adjustments[0] = 0.08
    except (AttributeError, IndexError, TypeError):
        pass
    tf = sh.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    _set_run_font(r, font_pt, color=CLR_TEXT)


def _arrow_label(slide, left, top, text: str = "→", font_pt: float = 11) -> None:
    tb = slide.shapes.add_textbox(left, top, Inches(0.22), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    _set_run_font(r, font_pt, bold=True, color=CLR_PRIMARY_SOFT)


def _add_pillar_card(slide, left, top, width, height, icon: str, title: str, body: str) -> None:
    """Карточка столпа (.bk-pillar): иконка, заголовок, текст."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = CLR_SURFACE
    card.line.color.rgb = CLR_BORDER
    card.line.width = Pt(0.75)
    try:
        card.adjustments[0] = 0.12
    except (AttributeError, IndexError, TypeError):
        pass
    tf = card.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.14)
    tf.margin_right = Inches(0.14)
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0.12)
    tf.vertical_anchor = MSO_ANCHOR.TOP

    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.LEFT
    ri = p0.add_run()
    ri.text = f"{icon}\n"
    _set_run_font(ri, 20, color=CLR_TEXT)

    p1 = tf.add_paragraph()
    p1.alignment = PP_ALIGN.LEFT
    rt = p1.add_run()
    rt.text = f"{title}\n"
    _set_run_font(rt, 13, bold=True, color=CLR_TEXT)

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    rb = p2.add_run()
    rb.text = body
    _set_run_font(rb, 10.5, color=CLR_TEXT_SOFT)


def _add_question_card(slide, left, top, width, height, question: str, domain: str) -> None:
    """Карточка вопроса (.bk-question-card)."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = CLR_SURFACE
    card.line.color.rgb = CLR_BORDER
    card.line.width = Pt(0.75)
    try:
        card.adjustments[0] = 0.1
    except (AttributeError, IndexError, TypeError):
        pass
    tf = card.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.08)
    tf.vertical_anchor = MSO_ANCHOR.TOP

    pq = tf.paragraphs[0]
    rq = pq.add_run()
    rq.text = f"“ {question}"
    _set_run_font(rq, 11, bold=False, color=CLR_TEXT)

    pd = tf.add_paragraph()
    rd = pd.add_run()
    rd.text = domain.upper()
    _set_run_font(rd, 7.5, bold=True, color=CLR_PRIMARY)
    try:
        rd.font.spacing = Pt(0.25)
    except (AttributeError, TypeError):
        pass


def _add_left_accent_panel(slide, left, top, width, height, heading: str, body: str) -> None:
    """Панель «Пример» / «Сценарий» (.bk-tech__example): белый фон, акцент слева."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = CLR_SURFACE
    card.line.color.rgb = CLR_BORDER
    card.line.width = Pt(0.75)
    try:
        card.adjustments[0] = 0.08
    except (AttributeError, IndexError, TypeError):
        pass
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.065), height)
    accent.fill.solid()
    accent.fill.fore_color.rgb = CLR_PRIMARY_SOFT
    try:
        accent.line.fill.background()
    except (AttributeError, TypeError):
        accent.line.color.rgb = CLR_PRIMARY_SOFT

    tf = card.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.22)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    tf.vertical_anchor = MSO_ANCHOR.TOP

    ph = tf.paragraphs[0]
    rh = ph.add_run()
    rh.text = f"{heading}\n"
    _set_run_font(rh, 11, bold=True, color=CLR_TEXT)

    pb = tf.add_paragraph()
    rb = pb.add_run()
    rb.text = body
    _set_run_font(rb, 10, color=CLR_TEXT_SOFT)


def _add_stack_group(slide, left, top, width, height, heading: str, lines: list[str]) -> None:
    """Колонка стека (.bk-stack__group)."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = CLR_SURFACE
    card.line.color.rgb = CLR_BORDER
    card.line.width = Pt(0.75)
    try:
        card.adjustments[0] = 0.1
    except (AttributeError, IndexError, TypeError):
        pass
    tf = card.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.08)
    tf.vertical_anchor = MSO_ANCHOR.TOP

    ph = tf.paragraphs[0]
    rh = ph.add_run()
    rh.text = f"{heading}\n\n"
    _set_run_font(rh, 11.5, bold=True, color=CLR_PRIMARY_DEEP)

    body_lines = "\n".join(f"• {ln}" for ln in lines)
    pb = tf.add_paragraph()
    rb = pb.add_run()
    rb.text = body_lines
    _set_run_font(rb, 10, color=CLR_TEXT)


def _add_problem_slide(prs: Presentation) -> None:
    """Слайд «Проблема» — как в Reveal: бейдж, заголовок, лид, три карточки статистики."""
    s = _blank_slide(prs)
    _add_eyebrow_pill(s, "Проблема", Inches(0.6), Inches(0.42), Inches(1.42), Inches(0.4))

    tb = s.shapes.add_textbox(Inches(0.6), Inches(0.95), Inches(12.0), Inches(0.85))
    tp = tb.text_frame.paragraphs[0]
    tr = tp.add_run()
    tr.text = "Знание есть — найти трудно"
    _set_run_font(tr, 28, bold=True, color=CLR_TEXT)

    lead = (
        "Регламенты в XWiki, инструкции в DOCX, FAQ в чатах, скрипты в общих папках. "
        "Сотрудник тратит время на поиск ответа, который у нас уже задокументирован."
    )
    lb = s.shapes.add_textbox(Inches(0.6), Inches(1.78), Inches(12.1), Inches(1.05))
    ltf = lb.text_frame
    ltf.word_wrap = True
    ltf.auto_size = MSO_AUTO_SIZE.NONE
    lp = ltf.paragraphs[0]
    lr = lp.add_run()
    lr.text = lead
    _set_run_font(lr, 16, color=CLR_TEXT)

    card_w = Inches(3.85)
    gap = Inches(0.25)
    top = Inches(3.05)
    h = Inches(2.35)
    stats = [
        ("5+", "источников документации\n(wiki, диск, почта, чаты)"),
        ("15–30 мин", "в среднем на поиск ответа\nпо типовому вопросу"),
        ("∞", "повторных вопросов\nв дежурные чаты поддержки"),
    ]
    x0 = Inches(0.55)
    for i, (val, lbl) in enumerate(stats):
        left = x0 + i * (card_w + gap)
        _add_stat_card(s, left, top, card_w, h, val, lbl)


def _add_solution_slide(prs: Presentation) -> None:
    """Слайд «Решение» — бейдж, три столпа с эмодзи-как в index.html."""
    s = _blank_slide(prs)
    _add_eyebrow_pill(s, "Решение", Inches(0.6), Inches(0.42), Inches(1.28), Inches(0.4))
    tb = s.shapes.add_textbox(Inches(0.6), Inches(0.95), Inches(12.0), Inches(0.55))
    tr = tb.text_frame.paragraphs[0].add_run()
    tr.text = "БочкарИИ — один чат для всей базы знаний"
    _set_run_font(tr, 28, bold=True, color=CLR_TEXT)
    lead = (
        "Сотрудник задаёт вопрос человеческим языком — система находит документы, "
        "показывает цитаты с релевантностью и формирует ответ. Всё локально."
    )
    lb = s.shapes.add_textbox(Inches(0.6), Inches(1.58), Inches(12.1), Inches(0.75))
    ltf = lb.text_frame
    ltf.word_wrap = True
    ltf.auto_size = MSO_AUTO_SIZE.NONE
    lr = ltf.paragraphs[0].add_run()
    lr.text = lead
    _set_run_font(lr, 15, color=CLR_TEXT)

    pillars = [
        ("💬", "Естественный язык", "Спрашиваем как человека: «как добавить пользователя в 1С?» — без ключевых слов и фильтров."),
        ("📑", "Ответы со ссылками", "Каждый ответ — с источниками, цитатами и оценкой релевантности. Можно перейти прямо в документ."),
        ("🔒", "Локально и безопасно", "Векторная база, LLM, документы — всё на наших серверах. Ничего не уходит в облако."),
    ]
    card_w = Inches(3.85)
    gap = Inches(0.28)
    top = Inches(2.72)
    h = Inches(3.95)
    x0 = Inches(0.55)
    for i, (ic, title, body) in enumerate(pillars):
        _add_pillar_card(s, x0 + i * (card_w + gap), top, card_w, h, ic, title, body)


def _add_questions_slide(prs: Presentation) -> None:
    """Шесть карточек вопросов в сетке 2×3."""
    s = _blank_slide(prs)
    _add_eyebrow_pill(s, "Что спрашивают на практике", Inches(0.6), Inches(0.42), Inches(3.55), Inches(0.4))
    tb = s.shapes.add_textbox(Inches(0.6), Inches(0.95), Inches(12.0), Inches(0.55))
    tr = tb.text_frame.paragraphs[0].add_run()
    tr.text = "Шесть запросов, которые ассистент закрывает каждый день"
    _set_run_font(tr, 26, bold=True, color=CLR_TEXT)

    questions = [
        ("Как сделать рассылку должникам?", "Маркетинг · CRM"),
        ("Как настроить состояние сертификации?", "Качество · процессы"),
        ("Инструкция по добавлению пользователя в 1С", "1С · администрирование"),
        ("Как найти ошибку в журнале регистрации?", "Эксплуатация · 1С"),
        ("Что делать с ошибкой «недостаточно прав для работы с таблицей»?", "1С · права доступа"),
        ("Количество поэкземплярной продукции не соответствует количеству подобранного товара", "Склад · ТСД"),
    ]
    cw = Inches(5.82)
    ch = Inches(1.02)
    x0, y0 = Inches(0.55), Inches(2.42)
    col_gap = Inches(0.28)
    row_gap = Inches(0.14)
    for i, (q, dom) in enumerate(questions):
        row, col = divmod(i, 2)
        left = x0 + col * (cw + col_gap)
        top = y0 + row * (ch + row_gap)
        _add_question_card(s, left, top, cw, ch, q, dom)

    foot = s.shapes.add_textbox(Inches(0.6), Inches(6.42), Inches(12.0), Inches(0.45))
    fr = foot.text_frame.paragraphs[0].add_run()
    fr.text = "Разные домены — 1С, склад, маркетинг, эксплуатация. Одна точка входа."
    _set_run_font(fr, 12.5, color=CLR_MUTED)


def _add_architecture_slide(prs: Presentation) -> None:
    """Блок-схема в рамке + краткий текст (аналог Mermaid)."""
    s = _blank_slide(prs)
    _add_eyebrow_pill(s, "Архитектура одним взглядом", Inches(0.6), Inches(0.42), Inches(2.65), Inches(0.4))
    tb = s.shapes.add_textbox(Inches(0.6), Inches(0.95), Inches(12.0), Inches(0.55))
    tr = tb.text_frame.paragraphs[0].add_run()
    tr.text = "Как вопрос превращается в ответ"
    _set_run_font(tr, 28, bold=True, color=CLR_TEXT)

    panel_top = Inches(1.56)
    panel_h = Inches(2.68)
    panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.48), panel_top, Inches(12.35), panel_h)
    panel.fill.solid()
    panel.fill.fore_color.rgb = CLR_SURFACE
    panel.line.color.rgb = CLR_BORDER
    panel.line.width = Pt(1)
    try:
        panel.adjustments[0] = 0.03
    except (AttributeError, IndexError, TypeError):
        pass

    y1 = Inches(1.68)
    h_node = Inches(0.4)
    row1: list[tuple] = [
        (Inches(0.86), "Сотрудник"),
        (Inches(1.0), "Web UI\nFlask + JS"),
        (Inches(0.98), "RAG ядро\ncore/rag.py"),
        (Inches(1.18), "Расширение\nrewrite, HyDE"),
        (Inches(1.05), "Гибридный\nпоиск"),
    ]
    x = Inches(0.62)
    hybrid_left = Inches(0)
    hybrid_w = Inches(0)
    gap_aw = Inches(0.12)
    for i, (w, txt) in enumerate(row1):
        if i == len(row1) - 1:
            hybrid_left = x
            hybrid_w = w
        _flow_node(s, x, y1, w, h_node, txt, 7.2)
        x = x + w
        if i < len(row1) - 1:
            _arrow_label(s, x, y1 + Inches(0.05))
            x = x + gap_aw

    cx = hybrid_left + hybrid_w / 2
    down = s.shapes.add_textbox(cx - Inches(0.08), Inches(2.08), Inches(0.2), Inches(0.28))
    dr = down.text_frame.paragraphs[0].add_run()
    dr.text = "↓"
    _set_run_font(dr, 12, bold=True, color=CLR_PRIMARY_SOFT)

    y2 = Inches(2.18)
    chroma_l = cx - Inches(1.18) - Inches(0.06)
    bm_l = cx + Inches(0.06)
    _flow_node(s, chroma_l, y2, Inches(1.12), h_node, "ChromaDB\nвекторы", 7.2)
    _flow_node(s, bm_l, y2, Inches(1.05), h_node, "BM25\nлексика", 7.2)

    dn2 = s.shapes.add_textbox(cx - Inches(0.08), Inches(2.58), Inches(0.2), Inches(0.26))
    d2 = dn2.text_frame.paragraphs[0].add_run()
    d2.text = "↓"
    _set_run_font(d2, 12, bold=True, color=CLR_PRIMARY_SOFT)

    y3 = Inches(2.68)
    _flow_node(s, cx - Inches(0.78), y3, Inches(1.56), Inches(0.38), "RRF +\ncross-encoder", 7.0)

    y4 = Inches(3.12)
    _flow_node(s, cx - Inches(0.72), y4, Inches(1.44), Inches(0.38), "LLM\nOllama / LM Studio", 7.0)

    y5 = Inches(3.52)
    _flow_node(s, cx - Inches(0.78), y5, Inches(1.56), Inches(0.34), "Ответ + источники\nв Web UI", 7.0)

    up = s.shapes.add_textbox(cx - Inches(0.35), Inches(3.88), Inches(0.7), Inches(0.22))
    ur = up.text_frame.paragraphs[0].add_run()
    ur.text = "↑ к пользователю"
    _set_run_font(ur, 7, color=CLR_MUTED)

    yk = Inches(4.02)
    _flow_node(s, Inches(1.85), yk, Inches(9.55), Inches(0.36), "База знаний (XWiki, PDF, DOCX, XLSX, HTML)  →  индексация в оба поиска", 7.0)

    lead_top = panel_top + panel_h + Inches(0.12)
    lead_txt = (
        "Коротко: вы отправляете текст вопроса — это не «разговор с интернетом». Система ищет ответы "
        "в ваших документах, собирает контекст и просит LLM сформулировать ответ."
    )
    lb = s.shapes.add_textbox(Inches(0.6), lead_top, Inches(12.1), Inches(0.55))
    lr = lb.text_frame.paragraphs[0].add_run()
    lr.text = lead_txt
    _set_run_font(lr, 9.5, color=CLR_TEXT_SOFT)

    steps_left = (
        "1. Вопрос из браузера (учёт истории при уточнениях).\n"
        "2. Подготовка к поиску: перефразирование, доп. формулировки.\n"
        "3. Параллельно векторный и BM25-поиск; RRF; при необходимости rerank."
    )
    steps_right = (
        "4. В модель — отобранные фрагменты из базы знаний.\n"
        "5. Локальная LLM формирует ответ по контексту.\n"
        "6. На экране — ответ, источники и цитаты для проверки."
    )
    sl = s.shapes.add_textbox(Inches(0.6), lead_top + Inches(0.52), Inches(5.95), Inches(1.05))
    r1 = sl.text_frame.paragraphs[0].add_run()
    r1.text = steps_left
    _set_run_font(r1, 9, color=CLR_TEXT)

    sr = s.shapes.add_textbox(Inches(6.55), lead_top + Inches(0.52), Inches(6.1), Inches(1.05))
    r2 = sr.text_frame.paragraphs[0].add_run()
    r2.text = steps_right
    _set_run_font(r2, 9, color=CLR_TEXT)


def _add_compact_flow_bar(s, left, top, width, height, text: str, font_pt: float = 7.5) -> None:
    """Одна полоса с цепочкой шагов в стиле .bk-flow--compact."""
    bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = CLR_SURFACE
    bar.line.color.rgb = CLR_BORDER
    try:
        bar.adjustments[0] = 0.35
    except (AttributeError, IndexError, TypeError):
        pass
    tf = bar.text_frame
    tf.clear()
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    _set_run_font(r, font_pt, color=CLR_TEXT)


def _add_hybrid_slide(prs: Presentation) -> None:
    s = _blank_slide(prs)
    _add_eyebrow_pill(s, "Под капотом · 1 / 3", Inches(0.6), Inches(0.4), Inches(1.85), Inches(0.38))
    tt = s.shapes.add_textbox(Inches(0.6), Inches(0.9), Inches(11.5), Inches(0.5))
    t0 = tt.text_frame.paragraphs[0].add_run()
    t0.text = "Гибридный поиск"
    _set_run_font(t0, 26, bold=True, color=CLR_TEXT)
    intro = (
        "Одной технологии поиска недостаточно: кто-то формулирует как в регламенте, кто-то — своими словами. "
        "Запрос идёт двумя дорожками, результаты склеиваются — не теряем ни точное слово, ни смысловой абзац."
    )
    ib = s.shapes.add_textbox(Inches(0.6), Inches(1.38), Inches(12.0), Inches(0.55))
    ir = ib.text_frame.paragraphs[0].add_run()
    ir.text = intro
    _set_run_font(ir, 10, color=CLR_TEXT_SOFT)

    flow_line = "Запрос   →   BM25   +   Dense   →   RRF   →   Rerank   →   Top-K"
    _add_compact_flow_bar(s, Inches(0.58), Inches(1.98), Inches(12.15), Inches(0.34), flow_line, 7.2)

    bullets = (
        "• По словам (BM25) — как умный Ctrl+F: «журнал регистрации», «ЭЦП», «ТСД».\n"
        "• По смыслу (Chroma) — другие формулировки ведут к тому же регламенту.\n"
        "• RRF — два топа кандидатов в одном рейтинге.\n"
        "• Rerank (опционально) — пересортировка по релевантности к вопросу."
    )
    bl = s.shapes.add_textbox(Inches(0.6), Inches(2.42), Inches(12.0), Inches(0.88))
    br = bl.text_frame.paragraphs[0].add_run()
    br.text = bullets
    _set_run_font(br, 9.5, color=CLR_TEXT)

    ex1 = (
        "«Ошибка в журнале регистрации» — BM25 цепляет дословно; если в wiki написано "
        "«просмотр событий», помогает поиск по смыслу."
    )
    ex2 = (
        "«Рассылка должникам» vs «информирование о задолженности» в документе — смысл тянет нужный блок; "
        "BM25 подхватывает «должник», «рассылка», если они есть в тексте."
    )
    ex3 = "В промпт попадают лучшие отрывки из вашей базы после слияния и (если включено) rerank — не «интернет из головы», а ваши документы."

    _add_left_accent_panel(s, Inches(0.6), Inches(3.38), Inches(5.75), Inches(1.1), "Пример 1 — термины 1С", ex1)
    _add_left_accent_panel(s, Inches(6.5), Inches(3.38), Inches(6.0), Inches(1.1), "Пример 2 — «живой» вопрос", ex2)
    _add_left_accent_panel(s, Inches(0.6), Inches(4.58), Inches(11.9), Inches(0.72), "Итог для модели", ex3)

    fn = s.shapes.add_textbox(Inches(0.6), Inches(5.38), Inches(12.0), Inches(0.4))
    fr = fn.text_frame.paragraphs[0].add_run()
    fr.text = ".env: RETRIEVAL_MODE, RRF_K_CONSTANT, RERANK_ENABLED, RERANK_MODEL · core/retrieval.py"
    _set_run_font(fr, 7.5, color=CLR_MUTED)


def _add_chunking_slide(prs: Presentation) -> None:
    s = _blank_slide(prs)
    _add_eyebrow_pill(s, "Под капотом · 2 / 3", Inches(0.6), Inches(0.4), Inches(1.85), Inches(0.38))
    h2 = s.shapes.add_textbox(Inches(0.6), Inches(0.9), Inches(11.5), Inches(0.5))
    h2r = h2.text_frame.paragraphs[0].add_run()
    h2r.text = "Структурное чанкирование"
    _set_run_font(h2r, 26, bold=True, color=CLR_TEXT)
    intro = (
        "Документ нарезаем не «через каждые 500 символов», а по структуре: заголовки, списки, таблицы. "
        "К каждому чанку добавляем контекст раздела."
    )
    i0 = s.shapes.add_textbox(Inches(0.6), Inches(1.38), Inches(12.0), Inches(0.55))
    i0r = i0.text_frame.paragraphs[0].add_run()
    i0r.text = intro
    _set_run_font(i0r, 10, color=CLR_TEXT_SOFT)

    btxt = (
        "• HTML — обход DOM: h1–h6, ul/ol, table.\n"
        "• DOCX — стили Heading*, таблицы по строкам.\n"
        "• Метаданные: section_path, chunk_kind — путь раздела виден в источниках.\n"
        "• Contextual Retrieval: LLM добавляет к чанку короткую аннотацию-префикс перед эмбеддингом."
    )
    bb = s.shapes.add_textbox(Inches(0.6), Inches(1.95), Inches(12.0), Inches(0.95))
    bbr = bb.text_frame.paragraphs[0].add_run()
    bbr.text = btxt
    _set_run_font(bbr, 10, color=CLR_TEXT)

    ex_body = (
        "Регламент 1С на 40 страниц делится по разделам — ассистент находит именно пункт "
        "«Раздача прав на конкретный документ», а не возвращает весь документ или невнятный кусок из середины."
    )
    _add_left_accent_panel(s, Inches(0.6), Inches(2.95), Inches(12.0), Inches(1.1), "Пример", ex_body)


def _add_memory_slide(prs: Presentation) -> None:
    s = _blank_slide(prs)
    _add_eyebrow_pill(s, "Под капотом · 3 / 3", Inches(0.6), Inches(0.4), Inches(1.85), Inches(0.38))
    h2 = s.shapes.add_textbox(Inches(0.6), Inches(0.9), Inches(11.5), Inches(0.5))
    h2r = h2.text_frame.paragraphs[0].add_run()
    h2r.text = "Память диалога"
    _set_run_font(h2r, 26, bold=True, color=CLR_TEXT)
    intro = (
        "Уточняющие вопросы «а для админа?» работают, потому что мы переписываем запрос с учётом истории, "
        "добавляем альтернативные формулировки и гипотетический ответ (HyDE)."
    )
    i0 = s.shapes.add_textbox(Inches(0.6), Inches(1.38), Inches(12.0), Inches(0.55))
    i0r = i0.text_frame.paragraphs[0].add_run()
    i0r.text = intro
    _set_run_font(i0r, 10, color=CLR_TEXT_SOFT)

    btxt = (
        "• Rewrite — уточнение в один самодостаточный запрос для поиска.\n"
        "• Multi-query — 2–3 альтернативные формулировки, объединение результатов.\n"
        "• HyDE — гипотетический ответ как дополнительный dense-запрос.\n"
        "• Переписанный запрос в SQLite (retrieval_query_text) — отладка."
    )
    bb = s.shapes.add_textbox(Inches(0.6), Inches(1.95), Inches(12.0), Inches(0.85))
    bbr = bb.text_frame.paragraphs[0].add_run()
    bbr.text = btxt
    _set_run_font(bbr, 10, color=CLR_TEXT)

    sc_body = (
        "Пользователь: «Как добавить пользователя в 1С?» → ассистент отвечает → "
        "Пользователь: «А как для роли админа?» → система переписывает запрос в "
        "«Как добавить пользователя 1С с ролью администратор» и находит правильный раздел."
    )
    _add_left_accent_panel(s, Inches(0.6), Inches(2.85), Inches(12.0), Inches(1.25), "Сценарий", sc_body)


def _add_stack_slide(prs: Presentation) -> None:
    s = _blank_slide(prs)
    _add_eyebrow_pill(s, "Стек и интеграции", Inches(0.6), Inches(0.42), Inches(1.65), Inches(0.4))
    tb = s.shapes.add_textbox(Inches(0.6), Inches(0.95), Inches(12.0), Inches(0.5))
    tr = tb.text_frame.paragraphs[0].add_run()
    tr.text = "Из чего собрано"
    _set_run_font(tr, 28, bold=True, color=CLR_TEXT)

    groups = [
        ("Бэкенд", ["Python 3.10+", "Flask + Flask-Cors", "SQLite (история, пользователи)"]),
        ("Поиск", ["ChromaDB (dense)", "rank-bm25 (sparse)", "sentence-transformers (rerank)"]),
        ("LLM", ["Ollama (Docker / native)", "LM Studio", "bge-m3, qwen2.5"]),
        ("Источники / каналы", ["Импорт XWiki", "PDF, DOCX, XLSX, PPTX, HTML"]),
    ]
    gw = Inches(2.95)
    gh = Inches(3.35)
    gap = Inches(0.22)
    y = Inches(1.68)
    x0 = Inches(0.52)
    for i, (head, lines) in enumerate(groups):
        _add_stack_group(s, x0 + i * (gw + gap), y, gw, gh, head, lines)

    tail = s.shapes.add_textbox(Inches(0.6), Inches(5.25), Inches(12.0), Inches(0.65))
    tr2 = tail.text_frame.paragraphs[0].add_run()
    tr2.text = (
        "Запуск: python web_app.py или start.bat. Документация: README.md, docs/user_guide.md, "
        "docs/CHANGELOG_RAG_2026-05.md."
    )
    _set_run_font(tr2, 9.5, color=CLR_MUTED)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. Титул
    _add_title_slide(prs)

    # 2. Проблема (бейдж + карточки, как в index.html / styles.css)
    _add_problem_slide(prs)

    # 3. Решение (столпы с иконками)
    _add_solution_slide(prs)

    # 4. Реальные вопросы (карточки)
    _add_questions_slide(prs)

    # 5. Архитектура (блок-схема + текст)
    _add_architecture_slide(prs)

    def add_demo_slide(eyebrow: str, title: str, image_name: str, bullets: list[str]):
        nonlocal prs
        slide = _blank_slide(prs)
        _section_header(slide, eyebrow, title)
        img_path = IMAGES / image_name
        if img_path.is_file():
            # Картинка слева ~58% ширины
            slide.shapes.add_picture(str(img_path), Inches(0.5), Inches(1.65), width=Inches(6.9))
        else:
            miss = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(6.9), Inches(1.0))
            mr = miss.text_frame.paragraphs[0].add_run()
            mr.text = f"(нет файла: {image_name})"
            _set_run_font(mr, 12, color=CLR_MUTED)

        tb = slide.shapes.add_textbox(Inches(7.65), Inches(1.65), Inches(5.0), Inches(5.5))
        tf = tb.text_frame
        tf.word_wrap = True
        for j, b in enumerate(bullets):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.level = 0
            p.space_after = Pt(6)
            r = p.add_run()
            r.text = f"• {b}"
            _set_run_font(r, 13, color=CLR_TEXT)

    # 6–11 Демо
    add_demo_slide(
        "Демо · 1 / 6",
        "Вход и роли",
        "01-login.png",
        [
            "Логин по email или username, пароль через корпоративную учётку.",
            "Две роли: user — чат и история, admin — плюс база знаний и диагностика.",
            "Гостевой режим только для чтения — в сайдбаре знакомые ранее чаты.",
            "Сессия сохраняется — переоткрыли вкладку, продолжили работу.",
        ],
    )
    add_demo_slide(
        "Демо · 2 / 6",
        "Новый чат и вопрос",
        "02-new-chat.png",
        [
            "«Новый чат» создаётся в один клик — слева в сайдбаре.",
            "Пишем естественным языком: «Инструкция по добавлению пользователя в 1С».",
            "В тулбаре — стиль ответа, сколько источников искать, минимальная релевантность.",
            "История переписывается — уточняющий вопрос «а как для роли админа?» работает.",
        ],
    )
    add_demo_slide(
        "Демо · 3 / 6",
        "Ответ с источниками",
        "Screenshot_10.jpg",
        [
            "Под ответом — кнопка «Источники»: справа выезжает панель с цитатами.",
            "Видно релевантность каждого источника и путь до раздела документа.",
            "Кнопка «Открыть» ведёт в исходный файл.",
            "Оценка «Полезно / Не полезно» — обратная связь по качеству ответов.",
        ],
    )
    add_demo_slide(
        "Демо · 4 / 6",
        "История диалогов",
        "04-chat-list.png",
        [
            "Все чаты сохраняются автоматически с понятными заголовками.",
            "Поиск по истории — быстро найти прошлый диалог.",
            "Переименование, удаление одного чата или очистка всей истории.",
            "«Скачать диалог» — выгрузка в текстовый файл.",
        ],
    )
    add_demo_slide(
        "Демо · 5 / 6 · только admin",
        "База знаний",
        "Screenshot_11.jpg",
        [
            "PDF, DOCX, XLSX, PPTX, HTML — нарезка и индексация.",
            "Предпросмотр показывает разбиение на чанки до индексации.",
            "«Переиндексировать» — фоновая задача, статус в UI.",
            "Актуализация регламентов без участия разработки.",
        ],
    )
    add_demo_slide(
        "Демо · 6 / 6 · только admin",
        "Админ-консоль",
        "Screenshot_12.jpg",
        [
            "Статус LLM, размер коллекции ChromaDB, активные модели.",
            "Метрики: диалоги, сообщения, оценки.",
            "Риски качества: устаревшие документы, дубли, слабая релевантность.",
            "Всё нужное дежурному админу на одной странице.",
        ],
    )

    # 12. Гибридный поиск + примеры
    _add_hybrid_slide(prs)

    # 13. Чанкинг + пример
    _add_chunking_slide(prs)

    # 14. Память диалога + сценарий
    _add_memory_slide(prs)

    # 15. Стек (карточки колонок)
    _add_stack_slide(prs)

    # 16. Q&A
    s = _blank_slide(prs)
    qa = s.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(4.0))
    tf = qa.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.CENTER
    r0 = p0.add_run()
    r0.text = "СПАСИБО\n\n"
    _set_run_font(r0, 14, bold=True, color=CLR_PRIMARY)
    r1 = p0.add_run()
    r1.text = "Вопросы?\n\n"
    _set_run_font(r1, 36, bold=True, color=CLR_TEXT)
    r2 = p0.add_run()
    r2.text = "Попробовать — http://localhost:5000\nИнструкция — docs/user_guide.md\n\n"
    _set_run_font(r2, 14, color=CLR_MUTED)
    r3 = p0.add_run()
    r3.text = (
        "Один вопрос в чат — и ответ из ваших регламентов и инструкций, "
        "без полчаса поиска по папкам и переписке."
    )
    _set_run_font(r3, 13, color=CLR_TEXT)

    prs.save(str(OUT))
    print(f"Сохранено: {OUT}", file=sys.stderr)


if __name__ == "__main__":
    build()
