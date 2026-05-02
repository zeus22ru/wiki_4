#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Скрипт для создания векторной базы данных из файлов в папке data/
Поддерживает HTML, DOCX, PDF, XLSX, XLS, PPTX, DOC форматы.
Использует ollama для генерации эмбеддингов и ChromaDB для хранения.
"""

import re
import sys
import requests
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path
from bs4 import BeautifulSoup
import chromadb
from chromadb.config import Settings
import json
from typing import List, Dict, Optional, Callable
import hashlib

# Импорт конфигурации и логирования
from config import settings, get_logger, inference_server_reachable, fetch_remote_model_ids

# Импорт общих функций для работы с эмбеддингами
from utils.embeddings import get_embedding, get_embeddings_batch, invalidate_embedding_cache

# Получаем логгер для этого модуля
logger = get_logger(__name__)

# Библиотеки для обработки разных форматов
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import pdfplumber
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from openpyxl import load_workbook
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False

try:
    import xlrd
    XLS_AVAILABLE = True
except ImportError:
    XLS_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import docx2txt
    DOC_AVAILABLE = True
except ImportError:
    DOC_AVAILABLE = False

# Устанавливаем UTF-8 для вывода в консоль (Windows) только при прямом запуске скрипта.
if sys.platform == 'win32' and __name__ == '__main__':
    import io
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
    sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8', errors='replace')

# Конфигурация загружается из config/settings.py
# OLLAMA_URL, OLLAMA_MODEL, OLLAMA_CHAT_MODEL, CHROMA_PERSIST_DIR,
# DATA_DIR, CHUNK_SIZE, CHUNK_OVERLAP, BATCH_SIZE


def _relative_source_path(file_path: Path) -> str:
    """Вернуть путь относительно DATA_DIR, а для временных файлов — имя файла."""
    try:
        return str(file_path.relative_to(settings.DATA_DIR))
    except ValueError:
        return file_path.name


def extract_text_from_html(html_path: Path) -> Optional[Dict[str, str]]:
    """Извлечь текст и метаданные из HTML файла"""
    try:
        with open(html_path, 'r', encoding='utf-8') as f:
            html_content = f.read()
        
        soup = BeautifulSoup(html_content, 'html.parser')
        
        # Удаляем скрипты и стили
        for script in soup(["script", "style", "nav", "footer", "header"]):
            script.decompose()
        
        # Получаем заголовок
        title = ""
        title_tag = soup.find('title')
        if title_tag:
            title = title_tag.get_text().strip()
        
        # Получаем h1
        h1 = ""
        h1_tag = soup.find('h1')
        if h1_tag:
            h1 = h1_tag.get_text().strip()
        
        # Получаем основной текст
        text = soup.get_text(separator=' ', strip=True)
        
        # Очистка текста
        text = re.sub(r'\s+', ' ', text)
        text = text.strip()
        
        return {
            "title": title or h1 or Path(html_path).stem,
            "content": text,
            "path": _relative_source_path(html_path)
        }
    except Exception as e:
        print(f"Ошибка при чтении {html_path}: {e}")
        return None


def extract_text_from_txt(txt_path: Path) -> Optional[Dict[str, str]]:
    """Извлечь текст из TXT файла."""
    try:
        text = txt_path.read_text(encoding='utf-8', errors='replace')
        text = re.sub(r'\s+', ' ', text).strip()
        return {
            "title": Path(txt_path).stem,
            "content": text,
            "path": _relative_source_path(txt_path),
        }
    except Exception as e:
        logger.error(f"Ошибка при чтении TXT {txt_path}: {e}")
        return None


def extract_text_from_docx(docx_path: Path) -> Optional[Dict[str, str]]:
    """Извлечь текст и метаданные из DOCX файла"""
    if not DOCX_AVAILABLE:
        logger.warning(f"Библиотека python-docx не установлена. Пропуск: {docx_path.name}")
        return None
    
    try:
        doc = Document(docx_path)
        
        # Извлекаем текст из параграфов
        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text.strip())
        
        # Извлекаем текст из таблиц
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    paragraphs.append(" | ".join(row_text))
        
        text = "\n".join(paragraphs)
        
        # Очистка текста
        text = re.sub(r'\s+', ' ', text)
        text = text.strip()
        
        # Получаем заголовок из свойств документа или первого параграфа
        title = doc.core_properties.title or ""
        if not title and paragraphs:
            title = paragraphs[0][:100]
        
        return {
            "title": title or Path(docx_path).stem,
            "content": text,
            "path": _relative_source_path(docx_path)
        }
    except Exception as e:
        logger.error(f"Ошибка при чтении DOCX {docx_path}: {e}")
        return None


def extract_text_from_pdf(pdf_path: Path) -> Optional[Dict[str, str]]:
    """Извлечь текст и метаданные из PDF файла"""
    if not PDF_AVAILABLE:
        logger.warning(f"Библиотека pdfplumber не установлена. Пропуск: {pdf_path.name}")
        return None
    
    try:
        text_parts = []
        title = ""
        
        with pdfplumber.open(pdf_path) as pdf:
            # Пытаемся получить заголовок из метаданных
            if pdf.metadata:
                title = pdf.metadata.get('Title', '') or pdf.metadata.get('Title', '')
            
            # Извлекаем текст со всех страниц
            for page_num, page in enumerate(pdf.pages, 1):
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)
                    
                    # Если заголовок не найден, берем первую строку первой страницы
                    if not title and page_num == 1:
                        lines = page_text.split('\n')
                        if lines:
                            title = lines[0].strip()[:100]
        
        text = "\n".join(text_parts)
        
        # Очистка текста
        text = re.sub(r'\s+', ' ', text)
        text = text.strip()
        
        return {
            "title": title or Path(pdf_path).stem,
            "content": text,
            "path": _relative_source_path(pdf_path)
        }
    except Exception as e:
        logger.error(f"Ошибка при чтении PDF {pdf_path}: {e}")
        return None


def extract_text_from_xlsx(xlsx_path: Path) -> Optional[Dict[str, str]]:
    """Извлечь текст и метаданные из XLSX файла"""
    if not XLSX_AVAILABLE:
        logger.warning(f"Библиотека openpyxl не установлена. Пропуск: {xlsx_path.name}")
        return None
    
    try:
        wb = load_workbook(xlsx_path, read_only=True, data_only=True)
        text_parts = []
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_text = []
            
            for row in sheet.iter_rows(values_only=True):
                row_values = [str(cell) if cell is not None else "" for cell in row]
                row_text = " | ".join(row_values).strip()
                if row_text:
                    sheet_text.append(row_text)
            
            if sheet_text:
                text_parts.append(f"Лист: {sheet_name}\n" + "\n".join(sheet_text))
        
        wb.close()
        
        text = "\n\n".join(text_parts)
        
        # Очистка текста
        text = re.sub(r'\s+', ' ', text)
        text = text.strip()
        
        return {
            "title": Path(xlsx_path).stem,
            "content": text,
            "path": _relative_source_path(xlsx_path)
        }
    except Exception as e:
        logger.error(f"Ошибка при чтении XLSX {xlsx_path}: {e}")
        return None


def extract_text_from_xls(xls_path: Path) -> Optional[Dict[str, str]]:
    """Извлечь текст и метаданные из XLS файла (старый формат Excel)"""
    if not XLS_AVAILABLE:
        logger.warning(f"Библиотека xlrd не установлена. Пропуск: {xls_path.name}")
        return None
    
    try:
        wb = xlrd.open_workbook(xls_path)
        text_parts = []
        
        for sheet_idx in range(wb.nsheets):
            sheet = wb.sheet_by_index(sheet_idx)
            sheet_text = []
            
            for row_idx in range(sheet.nrows):
                row_values = []
                for col_idx in range(sheet.ncols):
                    cell = sheet.cell_value(row_idx, col_idx)
                    if cell:
                        row_values.append(str(cell))
                
                if row_values:
                    sheet_text.append(" | ".join(row_values))
            
            if sheet_text:
                text_parts.append(f"Лист: {sheet.name}\n" + "\n".join(sheet_text))
        
        text = "\n\n".join(text_parts)
        
        # Очистка текста
        text = re.sub(r'\s+', ' ', text)
        text = text.strip()
        
        return {
            "title": Path(xls_path).stem,
            "content": text,
            "path": _relative_source_path(xls_path)
        }
    except Exception as e:
        logger.error(f"Ошибка при чтении XLS {xls_path}: {e}")
        return None


def extract_text_from_pptx(pptx_path: Path) -> Optional[Dict[str, str]]:
    """Извлечь текст и метаданные из PPTX файла"""
    if not PPTX_AVAILABLE:
        logger.warning(f"Библиотека python-pptx не установлена. Пропуск: {pptx_path.name}")
        return None
    
    try:
        prs = Presentation(pptx_path)
        text_parts = []
        title = ""
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            
            # Извлекаем текст со всех форм на слайде
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
                    
                    # Если заголовок не найден, берем текст с первого слайда
                    if not title and slide_num == 1:
                        title = shape.text.strip()[:100]
            
            if slide_text:
                text_parts.append(f"Слайд {slide_num}: " + " ".join(slide_text))
        
        text = "\n\n".join(text_parts)
        
        # Очистка текста
        text = re.sub(r'\s+', ' ', text)
        text = text.strip()
        
        return {
            "title": title or Path(pptx_path).stem,
            "content": text,
            "path": _relative_source_path(pptx_path)
        }
    except Exception as e:
        logger.error(f"Ошибка при чтении PPTX {pptx_path}: {e}")
        return None


def extract_text_from_doc(doc_path: Path) -> Optional[Dict[str, str]]:
    """Извлечь текст и метаданные из DOC файла (старый формат Word)"""
    if not DOC_AVAILABLE:
        logger.warning(f"Библиотека docx2txt не установлена. Пропуск: {doc_path.name}")
        return None
    
    try:
        text = docx2txt.process(doc_path)
        
        # Очистка текста
        text = re.sub(r'\s+', ' ', text)
        text = text.strip()
        
        # Берем первую строку как заголовок
        lines = text.split('\n')
        title = lines[0].strip()[:100] if lines else Path(doc_path).stem
        
        return {
            "title": title,
            "content": text,
            "path": _relative_source_path(doc_path)
        }
    except Exception as e:
        logger.error(f"Ошибка при чтении DOC {doc_path}: {e}")
        return None


def chunk_text(text: str, chunk_size: int = None, overlap: int = None) -> List[str]:
    """Разбить текст на чанки"""
    if chunk_size is None:
        chunk_size = settings.CHUNK_SIZE
    if overlap is None:
        overlap = settings.CHUNK_OVERLAP
    chunks = []
    start = 0
    text_len = len(text)
    
    while start < text_len:
        end = start + chunk_size
        chunk = text[start:end]
        
        # Пытаемся разбить по границе предложения
        if end < text_len:
            last_period = chunk.rfind('.')
            last_question = chunk.rfind('?')
            last_exclamation = chunk.rfind('!')
            last_boundary = max(last_period, last_question, last_exclamation)
            
            if last_boundary > chunk_size // 2:
                chunk = text[start:start + last_boundary + 1]
                end = start + last_boundary + 1
        
        chunks.append(chunk.strip())
        start = end - overlap
    
    return [c for c in chunks if len(c) > 50]


def get_file_handlers() -> Dict[str, Callable[[Path], Optional[Dict[str, str]]]]:
    """Поддерживаемые форматы файлов и функции извлечения текста."""
    return {
        '.html': extract_text_from_html,
        '.htm': extract_text_from_html,
        '.txt': extract_text_from_txt,
        '.docx': extract_text_from_docx,
        '.pdf': extract_text_from_pdf,
        '.xlsx': extract_text_from_xlsx,
        '.xls': extract_text_from_xls,
        '.pptx': extract_text_from_pptx,
        '.doc': extract_text_from_doc,
    }


def process_file(file_path: Path, extract_func: Callable[[Path], Optional[Dict[str, str]]]) -> List[Dict]:
    """Извлечь текст из одного файла и подготовить чанки для индексации."""
    ext = file_path.suffix.lower()
    doc_data = extract_func(file_path)

    if not doc_data or not doc_data["content"]:
        logger.warning(f"  Пропуск: не удалось извлечь текст из {file_path.name}")
        return []

    chunks = chunk_text(doc_data["content"])
    documents = []

    for j, chunk in enumerate(chunks):
        documents.append({
            "id": f"{hashlib.md5(f'{doc_data['path']}_{j}'.encode()).hexdigest()}",
            "text": chunk,
            "metadata": {
                "title": doc_data["title"],
                "source": doc_data["title"] or doc_data["path"],
                "path": doc_data["path"],
                "file_type": ext,
                "chunk_index": j,
                "total_chunks": len(chunks)
            }
        })

    return documents


def embed_documents_batch(batch_docs: List[Dict]) -> List[Dict]:
    """Получить эмбеддинги для пачки документов."""
    batch_texts = [doc["text"] for doc in batch_docs]
    batch_embeddings = get_embeddings_batch(batch_texts)
    embedded_documents = []

    if batch_embeddings and len(batch_embeddings) == len(batch_docs):
        for doc, embedding in zip(batch_docs, batch_embeddings):
            embedded_documents.append({
                "id": doc["id"],
                "text": doc["text"],
                "metadata": doc["metadata"],
                "embedding": embedding,
            })
        return embedded_documents

    logger.warning("Пакетная обработка эмбеддингов не удалась, пробуем по одному...")
    for doc in batch_docs:
        embedding = get_embedding(doc["text"])
        if embedding:
            embedded_documents.append({
                "id": doc["id"],
                "text": doc["text"],
                "metadata": doc["metadata"],
                "embedding": embedding,
            })
        else:
            logger.error(f"Не удалось получить эмбеддинг для документа {doc['id']}")

    return embedded_documents


def preview_document(file_path: Path, duplicate_exists: bool = False) -> Dict:
    """Проанализировать документ без записи в ChromaDB."""
    path = Path(file_path)
    ext = path.suffix.lower()
    handlers = get_file_handlers()
    warnings = []

    if ext not in handlers:
        return {
            "filename": path.name,
            "file_type": ext.lstrip("."),
            "size_bytes": path.stat().st_size if path.exists() else 0,
            "supported": False,
            "warnings": ["Формат файла не поддерживается"],
            "chunk_count": 0,
            "chunks": [],
            "title": path.stem,
            "text_length": 0,
        }

    if duplicate_exists:
        warnings.append("Файл с таким именем уже есть в базе знаний")

    size_bytes = path.stat().st_size
    if size_bytes > 50 * 1024 * 1024:
        warnings.append("Файл больше 50 MB, индексация может занять много времени")

    doc_data = handlers[ext](path)
    if not doc_data or not doc_data.get("content"):
        return {
            "filename": path.name,
            "file_type": ext.lstrip("."),
            "size_bytes": size_bytes,
            "supported": True,
            "warnings": warnings + ["Не удалось извлечь текст из файла"],
            "chunk_count": 0,
            "chunks": [],
            "title": path.stem,
            "text_length": 0,
        }

    content = doc_data["content"]
    chunks = chunk_text(content)
    if len(content) < 200:
        warnings.append("В документе мало извлеченного текста")
    if not chunks:
        warnings.append("После разбиения не получилось полезных чанков")

    headings = re.findall(r"(?:(?:^|[.!?])\s*)([А-ЯA-Z][^.!?]{8,80})", content)
    return {
        "filename": path.name,
        "file_type": ext.lstrip("."),
        "size_bytes": size_bytes,
        "supported": True,
        "warnings": warnings,
        "chunk_count": len(chunks),
        "chunks": chunks[:3],
        "title": doc_data.get("title") or path.stem,
        "text_length": len(content),
        "headings": headings[:5],
    }


def process_all_files(data_dir: str) -> List[Dict]:
    """Обработать все поддерживаемые файлы в директории"""
    data_path = Path(data_dir)
    documents = []
    
    logger.info(f"Сканирование директории: {data_path}")
    
    file_handlers = get_file_handlers()
    
    # Собираем все файлы поддерживаемых форматов
    all_files = []
    for ext in file_handlers.keys():
        files = list(data_path.rglob(f"*{ext}"))
        all_files.extend(files)
    
    # Исключаем временные файлы и файлы с расширениями, которые не нужно обрабатывать
    excluded_extensions = {'.crdownload', '.tmp', '.temp', '.bak'}
    all_files = [f for f in all_files if f.suffix.lower() not in excluded_extensions]
    
    logger.info(f"Найдено файлов: {len(all_files)}")
    
    # Группируем файлы по типу для статистики
    file_counts = {}
    for file_path in all_files:
        ext = file_path.suffix.lower()
        file_counts[ext] = file_counts.get(ext, 0) + 1
    
    logger.info("Статистика по типам файлов:")
    for ext, count in sorted(file_counts.items()):
        logger.info(f"  {ext}: {count}")
    
    worker_count = max(1, settings.DOCUMENT_PROCESS_WORKERS)
    if worker_count == 1 or len(all_files) <= 1:
        for i, file_path in enumerate(all_files, 1):
            ext = file_path.suffix.lower()

            if ext not in file_handlers:
                continue

            logger.info(f"Обработка {i}/{len(all_files)}: {file_path.name} ({ext})")
            documents.extend(process_file(file_path, file_handlers[ext]))
    else:
        logger.info(f"Параллельная обработка файлов: {worker_count} поток(ов)")
        results: List[List[Dict]] = [[] for _ in all_files]

        with ThreadPoolExecutor(max_workers=worker_count) as executor:
            futures = {}
            for i, file_path in enumerate(all_files):
                ext = file_path.suffix.lower()

                if ext not in file_handlers:
                    continue

                logger.info(f"Постановка в очередь {i + 1}/{len(all_files)}: {file_path.name} ({ext})")
                future = executor.submit(process_file, file_path, file_handlers[ext])
                futures[future] = (i, file_path)

            completed = 0
            for future in as_completed(futures):
                index, file_path = futures[future]
                completed += 1

                try:
                    results[index] = future.result()
                    logger.info(
                        f"Готово {completed}/{len(futures)}: {file_path.name}, чанков: {len(results[index])}"
                    )
                except Exception as e:
                    logger.error(f"Ошибка при обработке {file_path}: {e}")

        for file_documents in results:
            documents.extend(file_documents)
    
    logger.info(f"Всего создано чанков: {len(documents)}")
    return documents


def create_vector_db(documents: List[Dict], progress_callback: Optional[Callable[[Dict], None]] = None):
    """Создать векторную базу данных в ChromaDB с пакетной обработкой для GPU"""
    logger.info("Создание векторной базы данных...")

    def report_progress(progress: int, stage: str, message: str) -> None:
        if progress_callback:
            progress_callback({
                "progress": max(0, min(100, progress)),
                "stage": stage,
                "message": message,
            })
    
    # Создаем клиент ChromaDB. Старую коллекцию не трогаем, пока новые эмбеддинги не готовы.
    client = chromadb.PersistentClient(path=settings.CHROMA_PERSIST_DIR)
    report_progress(12, "prepare", "Подготовка векторной коллекции")
    
    # Подготавливаем данные для вставки
    ids = []
    texts = []
    metadatas = []
    embeddings = []
    
    total_docs = len(documents)
    batches = []
    for start in range(0, total_docs, settings.BATCH_SIZE):
        batch_end = min(start + settings.BATCH_SIZE, total_docs)
        batches.append((start, batch_end, documents[start:batch_end]))
    report_progress(15, "embedding", f"Генерация эмбеддингов: 0/{total_docs}")

    embedding_workers = max(1, settings.EMBEDDING_WORKERS)
    embedded_batches: List[List[Dict]] = [[] for _ in batches]

    if embedding_workers == 1 or len(batches) <= 1:
        embedded_count = 0
        for batch_index, (start, batch_end, batch_docs) in enumerate(batches):
            logger.info(
                f"Генерация эмбеддингов {start + 1}-{batch_end}/{total_docs} "
                f"(пакет {len(batch_docs)} документов)"
            )
            embedded_batches[batch_index] = embed_documents_batch(batch_docs)
            embedded_count += len(batch_docs)
            embedding_progress = 15 + int((embedded_count / max(total_docs, 1)) * 70)
            report_progress(
                embedding_progress,
                "embedding",
                f"Генерация эмбеддингов: {embedded_count}/{total_docs}",
            )
    else:
        logger.info(f"Параллельная генерация эмбеддингов: {embedding_workers} поток(ов)")
        with ThreadPoolExecutor(max_workers=embedding_workers) as executor:
            futures = {}
            for batch_index, (start, batch_end, batch_docs) in enumerate(batches):
                logger.info(
                    f"Постановка эмбеддингов в очередь {start + 1}-{batch_end}/{total_docs} "
                    f"(пакет {len(batch_docs)} документов)"
                )
                future = executor.submit(embed_documents_batch, batch_docs)
                futures[future] = (batch_index, start, batch_end)

            completed = 0
            embedded_count = 0
            for future in as_completed(futures):
                batch_index, start, batch_end = futures[future]
                completed += 1

                try:
                    embedded_batches[batch_index] = future.result()
                    embedded_count += batch_end - start
                    logger.info(
                        f"Эмбеддинги готовы {completed}/{len(futures)}: "
                        f"{start + 1}-{batch_end}/{total_docs}, "
                        f"получено: {len(embedded_batches[batch_index])}"
                    )
                    embedding_progress = 15 + int((embedded_count / max(total_docs, 1)) * 70)
                    report_progress(
                        embedding_progress,
                        "embedding",
                        f"Генерация эмбеддингов: {embedded_count}/{total_docs}",
                    )
                except Exception as e:
                    logger.error(f"Ошибка при генерации эмбеддингов {start + 1}-{batch_end}: {e}")

    for embedded_batch in embedded_batches:
        for doc in embedded_batch:
            ids.append(doc["id"])
            texts.append(doc["text"])
            metadatas.append(doc["metadata"])
            embeddings.append(doc["embedding"])

    if not ids:
        raise RuntimeError("Не удалось получить эмбеддинги: старая векторная база не изменена")
    
    # Вставляем данные в ChromaDB пакетами (максимальный размер пакета 5461)
    logger.info("Сохранение в ChromaDB...")
    report_progress(86, "saving", f"Сохранение в ChromaDB: 0/{len(ids)}")

    # Переключаем коллекцию только после успешной генерации эмбеддингов.
    try:
        client.delete_collection(settings.CHROMA_COLLECTION_NAME)
    except Exception:
        pass

    collection = client.create_collection(
        name=settings.CHROMA_COLLECTION_NAME,
        metadata={"description": "База знаний из XWiki"}
    )

    MAX_BATCH_SIZE = 5000  # Оставляем запас от лимита 5461
    
    total_docs = len(ids)
    saved = 0
    
    while saved < total_docs:
        batch_end = min(saved + MAX_BATCH_SIZE, total_docs)
        
        logger.info(f"Сохранение {saved+1}-{batch_end}/{total_docs} документов...")
        
        collection.add(
            ids=ids[saved:batch_end],
            documents=texts[saved:batch_end],
            metadatas=metadatas[saved:batch_end],
            embeddings=embeddings[saved:batch_end]
        )
        
        saved = batch_end
        save_progress = 86 + int((saved / max(total_docs, 1)) * 13)
        report_progress(save_progress, "saving", f"Сохранение в ChromaDB: {saved}/{total_docs}")
    
    logger.info(f"Векторная база данных создана! Всего документов: {len(ids)}")
    logger.info(f"База сохранена в: {settings.CHROMA_PERSIST_DIR}")
    
    # Инвалидируем кэш эмбеддингов после обновления базы
    logger.info("Инвалидация кэша эмбеддингов...")
    report_progress(99, "cache", "Очистка кэша эмбеддингов")
    invalidate_embedding_cache()
    logger.info("Кэш эмбеддингов очищен")


def main():
    """Главная функция"""
    logger.info("=" * 60)
    logger.info("Создание векторной базы знаний")
    logger.info("=" * 60)
    
    if not inference_server_reachable():
        logger.error(f"Сервер инференса недоступен: {settings.OLLAMA_URL}")
        logger.error(
            "Проверьте INFERENCE_BACKEND (ollama | lmstudio), запуск Ollama или LM Studio и загрузку моделей."
        )
        return
    logger.info(f"Сервер инференса отвечает: {settings.OLLAMA_URL}")

    try:
        model_names = fetch_remote_model_ids()
    except Exception as e:
        logger.error(f"Не удалось получить список моделей: {e}")
        return

    model_found = False
    for name in model_names:
        if name == settings.OLLAMA_EMBEDDING_MODEL or name.startswith(
            settings.OLLAMA_EMBEDDING_MODEL + ":"
        ):
            model_found = True
            logger.info(f"Модель для эмбеддингов: {name} ✓")
            break

    if not model_found:
        logger.error(f"ВНИМАНИЕ: Модель {settings.OLLAMA_EMBEDDING_MODEL} не найдена в списке сервера!")
        logger.error(f"Доступные модели: {', '.join(model_names)}")
        if settings.INFERENCE_BACKEND == "ollama" or settings.EMBEDDING_API_MODE == "ollama":
            logger.error(f"Установите модель: ollama pull {settings.OLLAMA_EMBEDDING_MODEL}")
        else:
            logger.error("В LM Studio загрузите модель эмбеддингов с тем же id, что в OLLAMA_EMBEDDING_MODEL.")
        return
    
    # Обрабатываем все поддерживаемые файлы
    documents = process_all_files(settings.DATA_DIR)
    
    if not documents:
        logger.warning("Не найдено документов для обработки")
        return
    
    # Создаем векторную базу данных
    create_vector_db(documents)
    
    logger.info("=" * 60)
    logger.info("Готово!")
    logger.info("=" * 60)


if __name__ == "__main__":
    main()
